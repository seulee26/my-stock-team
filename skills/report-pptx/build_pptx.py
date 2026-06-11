# -*- coding: utf-8 -*-
"""
report-pptx : 종목 리서치 스펙(JSON) -> 디자인된 PPTX 리포트 렌더러  (v2 — 증권사 리서치 톤)

사용:
    py build_pptx.py <spec.json> <out.pptx>

설계 원칙
- 이 스크립트는 "렌더링만" 담당한다. md 해석/요약은 호출하는 쪽(Claude)이 하고,
  결과를 SKILL.md의 SPEC 스키마(JSON)로 넘긴다.
- 디자인 규칙(전부 코드 고정):
    · 포인트색 KB 옐로우 #FFBC00 — 면적은 아끼고 액센트(바·마커·차트선)로만 쓴다
    · 본문 그레이/화이트, 표지는 잉크 배경. 차분한 금융 리서치 톤
    · 표는 면 채우기 대신 헤어라인 괘선(상·하 굵은 선 + 행 사이 얇은 선) — 증권사 표 문법
    · 한글 폰트 '맑은 고딕' 단일 고정(latin/ea/cs 모두 지정 -> 글자깨짐 방지)
    · 모든 수치 옆 (출처/기준일). source 없는 수치 항목은 렌더하지 않는다
    · 매수/매도·목표가 단정 금지(스펙 단계에서 거름)
    · 표 잘림 방지: 행 상한 + "…외 N행" 축약 + 셀 텍스트 길이 제한

슬라이드 순서(항상 고정, 7장):
    1 표지(종목명·작성일)  2 종목 개요  3 재무 요약(표+최근 3개년)
    4 가격/추세(차트)      5 뉴스·심리  6 리스크   7 한 줄 종합
"""

import sys, json, datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 디자인 토큰 ----------
FONT      = "맑은 고딕"
KB_YELLOW = RGBColor(0xFF, 0xBC, 0x00)
INK       = RGBColor(0x1A, 0x1A, 0x1A)   # 제목·강조
BODY      = RGBColor(0x40, 0x40, 0x40)   # 본문
MUTED     = RGBColor(0x97, 0x97, 0x97)   # 출처·보조
FAINT     = RGBColor(0xBD, 0xBD, 0xBD)   # 페이지번호 등 최약
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
PANEL     = RGBColor(0xF7, 0xF7, 0xF5)   # 아주 옅은 패널
HAIR      = RGBColor(0xDC, 0xDC, 0xD8)   # 헤어라인
COVER_BG  = RGBColor(0x17, 0x17, 0x17)   # 표지 잉크
COVER_SUB = RGBColor(0xB9, 0xB9, 0xB9)   # 표지 보조 텍스트

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.85)
CONTENT_W = SLIDE_W - 2 * MARGIN
HEAD_RULE_Y = Inches(1.46)
CONTENT_Y = Inches(1.78)
FOOT_Y = SLIDE_H - Inches(0.52)

MAX_FIN_ROWS  = 7
MAX_ISSUES    = 5
MAX_RISKS     = 4
MAX_MONITOR   = 4
CELL_MAXCHARS = 42
TEXT_MAXCHARS = 160

SECTION_NAMES = {
    2: ("01", "COMPANY OVERVIEW", "종목 개요"),
    3: ("02", "FINANCIAL SUMMARY", "재무 요약 — 최근 3개년"),
    4: ("03", "PRICE & MOMENTUM", "가격 · 추세"),
    5: ("04", "NEWS & SENTIMENT", "뉴스 · 심리"),
    6: ("05", "RISK ASSESSMENT", "리스크 점검"),
    7: ("06", "CONCLUSION", "한 줄 종합"),
}


# ---------- 텍스트/도형 헬퍼 ----------
def _apply_font(run, name=FONT, size=None, bold=None, italic=None, color=None, spacing=None):
    """폰트 고정(latin/ea/cs) + 선택적 자간(spacing: 1/100pt)."""
    run.font.name = name
    if size is not None:  run.font.size = Pt(size)
    if bold is not None:  run.font.bold = bold
    if italic is not None: run.font.italic = italic
    if color is not None: run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    if spacing is not None:
        rPr.set("spc", str(int(spacing)))
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def _trunc(s, n):
    s = "" if s is None else str(s)
    return s if len(s) <= n else s[: n - 1] + "…"


def _add_rect(slide, x, y, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _hairline(slide, x, y, w, color=HAIR, weight=1.0):
    return _add_rect(slide, x, y, w, Pt(weight), color)


def _textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    return tb, tf


def _para(tf, text, size, color, bold=False, italic=False, first=False,
          space_after=6, align=PP_ALIGN.LEFT, spacing=None, line=None):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if line is not None:
        p.line_spacing = line
    r = p.add_run()
    r.text = _trunc(text, TEXT_MAXCHARS)
    _apply_font(r, size=size, color=color, bold=bold, italic=italic, spacing=spacing)
    return p


def _marked_para(tf, text, size, color, first=False, space_after=4, line=1.12,
                 marker="▪", marker_color=KB_YELLOW, bold=False):
    """옐로우 마커 + 본문 한 단락."""
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(space_after)
    p.line_spacing = line
    rm = p.add_run(); rm.text = marker + "  "
    _apply_font(rm, size=size * 0.8, color=marker_color, bold=True)
    rt = p.add_run(); rt.text = _trunc(text, TEXT_MAXCHARS)
    _apply_font(rt, size=size, color=color, bold=bold)
    return p


def _slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def _header(slide, page, spec):
    """콘텐츠 공통 헤더: 섹션번호+영문 키커(자간) / 국문 제목 / 우측 종목·기준일 / 헤어라인."""
    no, kicker_en, title = SECTION_NAMES[page]
    # 키커: 노란 번호 + 자간 들어간 영문
    tb, tf = _textbox(slide, MARGIN, Inches(0.56), CONTENT_W, Inches(0.3))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); r1.text = no
    _apply_font(r1, size=10.5, color=KB_YELLOW, bold=True, spacing=200)
    r2 = p.add_run(); r2.text = "   " + kicker_en
    _apply_font(r2, size=10.5, color=MUTED, bold=True, spacing=200)
    # 제목
    tb2, tf2 = _textbox(slide, MARGIN, Inches(0.84), CONTENT_W - Inches(3.2), Inches(0.55))
    _para(tf2, title, 21, INK, bold=True, first=True, space_after=0)
    # 우측 종목명·기준일
    tb3, tf3 = _textbox(slide, SLIDE_W - MARGIN - Inches(3.6), Inches(0.62), Inches(3.6), Inches(0.7))
    p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.RIGHT
    r = p3.add_run(); r.text = spec.get("stock_name", "")
    _apply_font(r, size=12, color=INK, bold=True)
    tick = spec.get("ticker")
    if tick:
        r2t = p3.add_run(); r2t.text = "  " + str(tick)
        _apply_font(r2t, size=10, color=MUTED)
    p4 = tf3.add_paragraph(); p4.alignment = PP_ALIGN.RIGHT; p4.space_after = Pt(0)
    r4 = p4.add_run(); r4.text = "기준일 " + str(spec.get("as_of", ""))
    _apply_font(r4, size=9, color=MUTED)
    # 구분선: 짧은 옐로우 + 헤어라인
    _add_rect(slide, MARGIN, HEAD_RULE_Y, Inches(0.55), Pt(2.4), KB_YELLOW)
    _hairline(slide, MARGIN + Inches(0.55), HEAD_RULE_Y + Pt(0.7), CONTENT_W - Inches(0.55))


def _footer(slide, page, total=7):
    _hairline(slide, MARGIN, FOOT_Y, CONTENT_W)
    tb, tf = _textbox(slide, MARGIN, FOOT_Y + Inches(0.07), CONTENT_W - Inches(1.0), Inches(0.3))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "본 자료는 학습·연구 목적의 분석이며 투자 권유가 아닙니다."
    _apply_font(r, size=8, color=MUTED)
    r2 = p.add_run(); r2.text = "    |    STOCK-TEAM RESEARCH"
    _apply_font(r2, size=8, color=FAINT, spacing=100)
    tb2, tf2 = _textbox(slide, SLIDE_W - MARGIN - Inches(0.9), FOOT_Y + Inches(0.07), Inches(0.9), Inches(0.3))
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r3 = p2.add_run(); r3.text = "%02d / %02d" % (page, total)
    _apply_font(r3, size=8, color=FAINT)


def _source_line(slide, x, y, w, text, align=PP_ALIGN.LEFT):
    if not text:
        return
    tb, tf = _textbox(slide, x, y, w, Inches(0.28))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = _trunc(text, 130)
    _apply_font(r, size=8.5, color=MUTED)


def _empty_section(slide, msg="해당 자료 없음 (근거 부족)"):
    tb, tf = _textbox(slide, MARGIN, Inches(3.2), CONTENT_W, Inches(1.0), anchor=MSO_ANCHOR.MIDDLE)
    _para(tf, msg, 14, MUTED, first=True, align=PP_ALIGN.CENTER)


# ---------- 표 괘선(증권사 스타일) ----------
def _cell_borders(cell, top=None, bottom=None):
    """top/bottom: (RGBColor, width_pt) 또는 None. 좌우는 항상 제거."""
    tcPr = cell._tc.get_or_add_tcPr()
    for e in ("L", "R", "T", "B"):
        el = tcPr.find(qn("a:ln" + e))
        if el is not None:
            tcPr.remove(el)
    idx = 0
    for e in ("L", "R", "T", "B"):
        ln = tcPr.makeelement(qn("a:ln" + e), {})
        spec_ = {"T": top, "B": bottom}.get(e)
        if spec_:
            color, wpt = spec_
            ln.set("w", str(int(wpt * 12700)))
            ln.set("cap", "flat")
            sf = ln.makeelement(qn("a:solidFill"), {})
            c = sf.makeelement(qn("a:srgbClr"), {"val": str(color)})
            sf.append(c)
            ln.append(sf)
        else:
            ln.append(ln.makeelement(qn("a:noFill"), {}))
        tcPr.insert(idx, ln)
        idx += 1


def _fill_cell(cell, text, fs, color, bold, bg, align):
    cell.fill.solid(); cell.fill.fore_color.rgb = bg
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Pt(6); cell.margin_right = Pt(6)
    cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = _trunc(text, CELL_MAXCHARS)
    _apply_font(r, size=fs, color=color, bold=bold)


# ---------- 차트 축 스타일 ----------
def _style_axes(chart, label_size=8.5):
    try:
        chart.font.size = Pt(label_size)
        chart.font.name = FONT
        chart.font.color.rgb = MUTED
    except Exception:
        pass
    for ax_name in ("value_axis", "category_axis"):
        try:
            ax = getattr(chart, ax_name)
            ax.format.line.color.rgb = HAIR
            ax.tick_labels.font.size = Pt(label_size)
            ax.tick_labels.font.name = FONT
            ax.tick_labels.font.color.rgb = MUTED
        except Exception:
            pass
    try:
        va = chart.value_axis
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEC)
        va.major_gridlines.format.line.width = Pt(0.5)
    except Exception:
        pass


def _thin_cat_labels(chart, skip):
    try:
        catAx = chart._chartSpace.find(".//" + qn("c:catAx"))
        if catAx is None:
            return
        for tag in ("c:tickLblSkip", "c:tickMarkSkip"):
            el = catAx.find(qn(tag))
            if el is None:
                el = catAx.makeelement(qn(tag), {})
                catAx.append(el)
            el.set("val", str(skip))
    except Exception:
        pass


# ================= 슬라이드 1: 표지 =================
def slide_cover(prs, spec):
    s = _slide(prs)
    _add_rect(s, 0, 0, SLIDE_W, SLIDE_H, COVER_BG)
    # 최상단 옐로우 스트립
    _add_rect(s, 0, 0, SLIDE_W, Inches(0.09), KB_YELLOW)
    # 상단 좌: 발행 주체, 우: 문서 유형
    tb0, tf0 = _textbox(s, MARGIN, Inches(0.55), Inches(6), Inches(0.3))
    _para(tf0, "STOCK-TEAM RESEARCH", 10.5, COVER_SUB, bold=True, first=True, spacing=300)
    tbr, tfr = _textbox(s, SLIDE_W - MARGIN - Inches(5), Inches(0.55), Inches(5), Inches(0.3))
    _para(tfr, "EQUITY RESEARCH  ·  SHORT-TERM MOMENTUM (1W)", 9, COVER_SUB, first=True,
          align=PP_ALIGN.RIGHT, spacing=200)
    _hairline(s, MARGIN, Inches(0.95), CONTENT_W, RGBColor(0x33, 0x33, 0x33))
    # 가드레일: 리포트 첫머리 고지
    tbn, tfn = _textbox(s, MARGIN, Inches(1.07), CONTENT_W, Inches(0.28))
    _para(tfn, "무료 공개 데이터 기반 학습용", 9.5, RGBColor(0x8A, 0x8A, 0x8A), first=True)

    # 본문: 티커 키커 + 종목명
    tb, tf = _textbox(s, MARGIN, Inches(2.35), CONTENT_W, Inches(2.3))
    ticker = spec.get("ticker")
    kick = ("KRX %s" % ticker) if ticker else "EQUITY"
    _para(tf, kick, 13, KB_YELLOW, bold=True, first=True, space_after=10, spacing=250)
    _para(tf, spec.get("stock_name", "종목"), 52, WHITE, bold=True, space_after=8)

    # 종합 의견 라인 (옐로우 바 + 텍스트)
    overall = spec.get("overall")
    if overall:
        oy = Inches(4.55)
        _add_rect(s, MARGIN, oy, Pt(3), Inches(0.62), KB_YELLOW)
        tb2, tf2 = _textbox(s, MARGIN + Inches(0.22), oy, Inches(10.5), Inches(0.62),
                            anchor=MSO_ANCHOR.MIDDLE)
        p = tf2.paragraphs[0]
        r1 = p.add_run(); r1.text = "종합 모멘텀 의견   "
        _apply_font(r1, size=11, color=COVER_SUB, bold=True)
        r2 = p.add_run(); r2.text = _trunc(overall, 60)
        _apply_font(r2, size=16, color=WHITE, bold=True)

    # 하단 메타 행
    _hairline(s, MARGIN, Inches(6.35), CONTENT_W, RGBColor(0x33, 0x33, 0x33))
    created = spec.get("created") or spec.get("as_of") or datetime.date.today().isoformat()
    tb3, tf3 = _textbox(s, MARGIN, Inches(6.5), CONTENT_W, Inches(0.3))
    p3 = tf3.paragraphs[0]
    meta_pairs = [("작성일", str(created)), ("기준일", str(spec.get("as_of", created)))]
    first = True
    for label, val in meta_pairs:
        r = p3.add_run(); r.text = ("" if first else "        ") + label + "  "
        _apply_font(r, size=9.5, color=RGBColor(0x77, 0x77, 0x77))
        rv = p3.add_run(); rv.text = val
        _apply_font(rv, size=9.5, color=COVER_SUB, bold=True)
        first = False
    tb4, tf4 = _textbox(s, MARGIN, Inches(6.92), CONTENT_W, Inches(0.3))
    _para(tf4, "본 자료는 학습·연구 목적의 분석이며 투자 권유가 아닙니다. 매수·매도·목표가 단정을 포함하지 않습니다.",
          8.5, RGBColor(0x6E, 0x6E, 0x6E), first=True)


# ================= 슬라이드 2: 종목 개요 =================
def slide_overview(prs, spec, page=2):
    s = _slide(prs)
    _header(s, page, spec)
    ov = spec.get("overview") or {}
    bullets = ov.get("bullets") or []
    if not bullets:
        _empty_section(s)
    else:
        tb, tf = _textbox(s, MARGIN, CONTENT_Y + Inches(0.25), CONTENT_W, Inches(4.4))
        first = True
        for b in bullets[:6]:
            txt = b.get("text") if isinstance(b, dict) else str(b)
            src = b.get("source") if isinstance(b, dict) else None
            _marked_para(tf, txt, 14.5, INK if first else BODY, first=first, space_after=3,
                         bold=first)
            first = False
            if src:
                p = tf.add_paragraph(); p.space_after = Pt(14)
                r = p.add_run(); r.text = "      " + _trunc(src, 110)
                _apply_font(r, size=8.5, color=MUTED)
    _footer(s, page)


# ================= 슬라이드 3: 재무 요약(표) =================
def slide_financials(prs, spec, page=3):
    s = _slide(prs)
    _header(s, page, spec)
    fin = spec.get("financials") or {}
    headers = fin.get("headers") or []
    rows = fin.get("rows") or []
    if not headers or not rows:
        _empty_section(s)
        _footer(s, page); return

    overflow = max(0, len(rows) - MAX_FIN_ROWS)
    rows = rows[:MAX_FIN_ROWS]
    ncol = len(headers)
    nrow = len(rows) + 1

    # 단위 라벨(표 우상단)
    if fin.get("unit"):
        _source_line(s, MARGIN, CONTENT_Y, CONTENT_W, fin["unit"], align=PP_ALIGN.RIGHT)

    tbl_y = CONTENT_Y + Inches(0.32)
    body_fs = 12 if len(rows) <= 5 else 11
    row_h = Inches(0.58) if len(rows) <= 5 else Inches(0.5)
    tbl_h = row_h * nrow

    gframe = s.shapes.add_table(nrow, ncol, MARGIN, tbl_y, CONTENT_W, tbl_h)
    table = gframe.table
    table.first_row = False
    table.horz_banding = False

    first_w = int(CONTENT_W * 0.28)
    rest_w = int((CONTENT_W - first_w) / (ncol - 1)) if ncol > 1 else CONTENT_W
    table.columns[0].width = Emu(first_w)
    for c in range(1, ncol):
        table.columns[c].width = Emu(rest_w)

    # 헤더: 흰 배경 + 상단 진한 선 + 하단 진한 선 (증권사 표 문법)
    for c in range(ncol):
        align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
        cell = table.cell(0, c)
        _fill_cell(cell, headers[c], body_fs, INK, True, WHITE, align)
        _cell_borders(cell, top=(INK, 1.6), bottom=(INK, 1.1))
    # 데이터: 흰 배경 + 행 사이 헤어라인, 마지막 행 하단 진한 선
    for i, row in enumerate(rows):
        last = (i == len(rows) - 1)
        for c in range(ncol):
            val = row[c] if c < len(row) else ""
            cell = table.cell(i + 1, c)
            if c == 0:
                _fill_cell(cell, val, body_fs, INK, True, WHITE, PP_ALIGN.LEFT)
            else:
                neg = str(val).strip().startswith("-")
                _fill_cell(cell, val, body_fs, BODY, False, WHITE, PP_ALIGN.RIGHT)
            _cell_borders(cell, bottom=(INK, 1.4) if last else (HAIR, 0.75))

    y_after = tbl_y + tbl_h + Inches(0.18)
    if overflow:
        _source_line(s, MARGIN, y_after, CONTENT_W, "…외 %d행 생략(슬라이드 가독성)" % overflow)
        y_after += Inches(0.24)
    _source_line(s, MARGIN, y_after, CONTENT_W, fin.get("source", ""))
    _footer(s, page)


# ================= 슬라이드 4: 가격/추세(차트) =================
def _maybe_fetch_series(ticker, as_of):
    if not ticker:
        return None
    try:
        import FinanceDataReader as fdr
        end = as_of or datetime.date.today().isoformat()
        end_d = datetime.date.fromisoformat(str(end)[:10])
        start_d = end_d - datetime.timedelta(days=190)
        df = fdr.DataReader(str(ticker), start_d.isoformat(), end_d.isoformat())
        if df is None or df.empty or "Close" not in df.columns:
            return None
        out = [{"date": d.strftime("%m/%d"), "close": float(v)}
               for d, v in df["Close"].items() if v == v]
        return out or None
    except Exception:
        return None


def slide_price(prs, spec, page=4):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    s = _slide(prs)
    _header(s, page, spec)
    price = spec.get("price") or {}

    series = price.get("price_series")
    if not series:
        series = _maybe_fetch_series(spec.get("ticker"), spec.get("as_of"))

    chart_x, chart_y = MARGIN, CONTENT_Y + Inches(0.42)
    panel_w = Inches(3.55)
    chart_w = CONTENT_W - panel_w - Inches(0.5)
    chart_h = Inches(3.85)
    drew_chart = False

    # 차트 캡션
    tbc, tfc = _textbox(s, chart_x, CONTENT_Y, chart_w, Inches(0.3))
    cap = "종가 추이 — 최근 6개월" if series else "기간별 변동률"
    _para(tfc, cap, 11, INK, bold=True, first=True)

    try:
        if series and len(series) >= 5:
            step = max(1, len(series) // 130)
            sampled = series[::step]
            cats = [d.get("date", "") for d in sampled]
            vals = [d.get("close", 0) for d in sampled]
            cd = CategoryChartData()
            cd.categories = cats
            cd.add_series("종가", vals)
            gf = s.shapes.add_chart(XL_CHART_TYPE.LINE, chart_x, chart_y, chart_w, chart_h, cd)
            ch = gf.chart
            ch.has_legend = False
            ch.has_title = False
            ser = ch.series[0]
            ser.format.line.color.rgb = KB_YELLOW
            ser.format.line.width = Pt(2.5)
            ser.smooth = False
            _style_axes(ch)
            try:
                ch.value_axis.tick_labels.number_format = "#,##0"
                ch.value_axis.tick_labels.number_format_is_linked = False
            except Exception:
                pass
            _thin_cat_labels(ch, max(1, len(cats) // 6))
            drew_chart = True
        elif price.get("change_rates"):
            cr = price["change_rates"]
            cd = CategoryChartData()
            cd.categories = list(cr.keys())
            cd.add_series("변동률(%)", [float(v) for v in cr.values()])
            gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_x, chart_y, chart_w, chart_h, cd)
            ch = gf.chart
            ch.has_legend = False
            ch.has_title = False
            ch.series[0].format.fill.solid()
            ch.series[0].format.fill.fore_color.rgb = KB_YELLOW
            try:
                ch.plots[0].gap_width = 120
                ch.plots[0].has_data_labels = True
                dl = ch.plots[0].data_labels
                dl.number_format = '0.0"%"'
                dl.number_format_is_linked = False
                dl.font.size = Pt(10); dl.font.name = FONT; dl.font.color.rgb = INK
            except Exception:
                pass
            _style_axes(ch)
            drew_chart = True
    except Exception:
        drew_chart = False

    if not drew_chart:
        tb, tf = _textbox(s, chart_x, Inches(3.4), chart_w, Inches(1.0), anchor=MSO_ANCHOR.MIDDLE)
        _para(tf, "차트용 시계열/변동률 데이터 없음 (근거 부족)", 13, MUTED, first=True,
              align=PP_ALIGN.CENTER)

    # 우측 KEY DATA 패널 (헤어라인 행 구조)
    metrics = price.get("metrics") or []
    if metrics:
        px = SLIDE_W - MARGIN - panel_w
        tbp, tfp = _textbox(s, px, CONTENT_Y, panel_w, Inches(0.3))
        _para(tfp, "KEY DATA", 10, MUTED, bold=True, first=True, spacing=250)
        _hairline(s, px, CONTENT_Y + Inches(0.32), panel_w, INK, 1.4)
        ry = CONTENT_Y + Inches(0.42)
        row_h = Inches(0.585)
        for m in metrics[:7]:
            label = m.get("label", "") if isinstance(m, dict) else ""
            value = m.get("value", "") if isinstance(m, dict) else str(m)
            tbr_, tfr_ = _textbox(s, px, ry, panel_w, row_h, anchor=MSO_ANCHOR.MIDDLE)
            p = tfr_.paragraphs[0]
            rl = p.add_run(); rl.text = _trunc(label, 14)
            _apply_font(rl, size=10, color=MUTED)
            # 우측 정렬 값: 별도 단락 불가 → 우측 텍스트박스
            tbv, tfv = _textbox(s, px, ry, panel_w, row_h, anchor=MSO_ANCHOR.MIDDLE)
            pv = tfv.paragraphs[0]; pv.alignment = PP_ALIGN.RIGHT
            rv = pv.add_run(); rv.text = _trunc(value, 20)
            _apply_font(rv, size=12.5, color=INK, bold=True)
            ry += row_h
            _hairline(s, px, ry, panel_w)

    # 코멘트 + 출처
    y = chart_y + chart_h + Inches(0.16)
    if price.get("comment"):
        tb, tf = _textbox(s, MARGIN, y, chart_w, Inches(0.6))
        _marked_para(tf, price["comment"], 11.5, BODY, first=True, line=1.15)
        y += Inches(0.46)
    _source_line(s, MARGIN, y, chart_w, price.get("source", ""))
    _footer(s, page)


# ================= 슬라이드 5: 뉴스·심리 =================
def slide_sentiment(prs, spec, page=5):
    s = _slide(prs)
    _header(s, page, spec)
    sen = spec.get("sentiment") or {}
    verdict = sen.get("verdict")
    issues = sen.get("issues") or []

    y = CONTENT_Y + Inches(0.1)
    if verdict:
        _add_rect(s, MARGIN, y, Pt(3), Inches(0.52), KB_YELLOW)
        tb, tf = _textbox(s, MARGIN + Inches(0.2), y, Inches(8.5), Inches(0.52),
                          anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = "시장 심리   "
        _apply_font(r1, size=10.5, color=MUTED, bold=True)
        r2 = p.add_run(); r2.text = _trunc(verdict, 45)
        _apply_font(r2, size=15, color=INK, bold=True)
        y += Inches(0.78)
        _hairline(s, MARGIN, y - Inches(0.12), CONTENT_W)

    if not issues:
        if not verdict:
            _empty_section(s)
    else:
        tb, tf = _textbox(s, MARGIN, y + Inches(0.08), CONTENT_W, FOOT_Y - y - Inches(0.25))
        first = True
        for it in issues[:MAX_ISSUES]:
            txt = it.get("text", "") if isinstance(it, dict) else str(it)
            src = it.get("source", "") if isinstance(it, dict) else ""
            dt = it.get("date", "") if isinstance(it, dict) else ""
            _marked_para(tf, txt, 13, BODY, first=first, space_after=2, line=1.15)
            first = False
            meta = "   ·   ".join([x for x in [dt, ("출처: " + src) if src else ""] if x])
            if meta:
                p = tf.add_paragraph(); p.space_after = Pt(12)
                r = p.add_run(); r.text = "     " + _trunc(meta, 110)
                _apply_font(r, size=8.5, color=MUTED)
    _footer(s, page)


# ================= 슬라이드 6: 리스크 =================
def slide_risk(prs, spec, page=6):
    s = _slide(prs)
    _header(s, page, spec)
    risk = spec.get("risk") or {}
    items = risk.get("items") or []
    monitoring = risk.get("monitoring") or []

    if not items and not monitoring:
        _empty_section(s)
        _footer(s, page); return

    if items:
        tb, tf = _textbox(s, MARGIN, CONTENT_Y + Inches(0.12), CONTENT_W, Inches(3.3))
        first = True
        for i, it in enumerate(items[:MAX_RISKS], 1):
            title = it.get("title", "") if isinstance(it, dict) else str(it)
            detail = it.get("detail", "") if isinstance(it, dict) else ""
            src = it.get("source", "") if isinstance(it, dict) else ""
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.space_after = Pt(2)
            rn = p.add_run(); rn.text = "%02d" % i
            _apply_font(rn, size=13.5, color=KB_YELLOW, bold=True, spacing=100)
            rt = p.add_run(); rt.text = "   " + _trunc(title, 40)
            _apply_font(rt, size=13.5, color=INK, bold=True)
            first = False
            if detail:
                pd = tf.add_paragraph(); pd.space_after = Pt(1); pd.line_spacing = 1.12
                rd = pd.add_run(); rd.text = "       " + _trunc(detail, 130)
                _apply_font(rd, size=11.5, color=BODY)
            if src:
                ps = tf.add_paragraph(); ps.space_after = Pt(11)
                rs = ps.add_run(); rs.text = "       " + _trunc(src, 100)
                _apply_font(rs, size=8.5, color=MUTED)

    if monitoring:
        my = Inches(5.15)
        _hairline(s, MARGIN, my, CONTENT_W, INK, 1.2)
        tb, tf = _textbox(s, MARGIN, my + Inches(0.12), CONTENT_W, Inches(1.5))
        _para(tf, "MONITORING TRIGGERS  ·  관찰 트리거", 9.5, MUTED, bold=True, first=True,
              space_after=6, spacing=200)
        for m in monitoring[:MAX_MONITOR]:
            _marked_para(tf, _trunc(m, 100), 11, BODY, space_after=3, marker="–",
                         marker_color=KB_YELLOW)
    _footer(s, page)


# ================= 슬라이드 7: 한 줄 종합 =================
def slide_conclusion(prs, spec, page=7):
    s = _slide(prs)
    _header(s, page, spec)
    con = spec.get("conclusion") or {}
    headline = con.get("headline")
    basis = con.get("basis") or []

    # 헤드라인: 큰 타이포 + 좌측 옐로우 바 (박스 대신 타이포 중심)
    py = CONTENT_Y + Inches(0.45)
    _add_rect(s, MARGIN, py + Inches(0.06), Pt(4), Inches(1.25), KB_YELLOW)
    tb, tf = _textbox(s, MARGIN + Inches(0.32), py, CONTENT_W - Inches(0.32), Inches(1.6))
    _para(tf, headline or "종합 의견 (근거 부족)", 19.5, INK, bold=True, first=True,
          space_after=0, line=1.25)

    if basis:
        by = py + Inches(1.85)
        _hairline(s, MARGIN, by, CONTENT_W)
        tb2, tf2 = _textbox(s, MARGIN, by + Inches(0.14), CONTENT_W, Inches(2.2))
        _para(tf2, "KEY RATIONALE  ·  핵심 근거", 9.5, MUTED, bold=True, first=True,
              space_after=8, spacing=200)
        for b in basis[:4]:
            _marked_para(tf2, _trunc(b, 130), 12.5, BODY, space_after=7, line=1.15)

    # 가드레일: 리포트 끝 — 데이터 출처·기준일 목록 (spec 전체에서 자동 수집)
    srcs = []
    def _collect(v):
        v = (v or "").strip()
        if v and v not in srcs:
            srcs.append(v)
    _collect((spec.get("financials") or {}).get("source"))
    _collect((spec.get("price") or {}).get("source"))
    for b in (spec.get("overview") or {}).get("bullets") or []:
        if isinstance(b, dict):
            _collect(b.get("source"))
    for it in (spec.get("sentiment") or {}).get("issues") or []:
        if isinstance(it, dict) and it.get("source"):
            _collect("출처: " + it["source"] + ((", " + it["date"]) if it.get("date") else ""))
    for it in (spec.get("risk") or {}).get("items") or []:
        if isinstance(it, dict):
            _collect(it.get("source"))
    if srcs:
        sy = FOOT_Y - Inches(0.78)
        _hairline(s, MARGIN, sy, CONTENT_W)
        tb4, tf4 = _textbox(s, MARGIN, sy + Inches(0.08), CONTENT_W, Inches(0.62))
        _para(tf4, "DATA SOURCES  ·  데이터 출처·기준일", 8.5, MUTED, bold=True, first=True,
              space_after=3, spacing=150)
        _para(tf4, "   ·   ".join(_trunc(x, 55) for x in srcs[:6]), 8, MUTED,
              space_after=0, line=1.15)
    _footer(s, page)


# ---------- 메인 ----------
def build(spec, out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_cover(prs, spec)
    slide_overview(prs, spec)
    slide_financials(prs, spec)
    slide_price(prs, spec)
    slide_sentiment(prs, spec)
    slide_risk(prs, spec)
    slide_conclusion(prs, spec)
    prs.save(out_path)
    return out_path


def main():
    if len(sys.argv) < 3:
        print("usage: py build_pptx.py <spec.json> <out.pptx>", file=sys.stderr)
        sys.exit(2)
    spec_path, out_path = sys.argv[1], sys.argv[2]
    with open(spec_path, "r", encoding="utf-8") as f:
        spec = json.load(f)
    if not spec.get("created"):
        spec["created"] = datetime.date.today().isoformat()
    build(spec, out_path)
    print("OK ->", out_path)


if __name__ == "__main__":
    main()
